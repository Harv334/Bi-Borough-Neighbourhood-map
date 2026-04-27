#!/usr/bin/env python3
"""Generate a per-ward slide deck (.pptx) from ward_data.json + lsoa_data.json.

Mirrors the A4 ward profile in the dashboard but as a 5-slide PowerPoint deck:
  Slide 1 — Cover: name, borough, neighbourhood, KPIs
  Slide 2 — Demographics: ethnicity stack, age stack, vs-NWL compare bars
  Slide 3 — Health & deprivation: IMD radar + health bars
  Slide 4 — Crime: 11-category bar with NWL median ticks
  Slide 5 — Civic strength + green/blue space access

Usage:
    python scripts/build_ward_pptx.py --ward "St James's"
    python scripts/build_ward_pptx.py --ward "Roundwood" --out output/
    python scripts/build_ward_pptx.py --all                # bulk: every ward
    python scripts/build_ward_pptx.py --list               # print available

Output: nwl-ward-{slug}.pptx in --out (default: ./output/ward_decks/)

Dependencies:
    pip install python-pptx matplotlib --break-system-packages
"""
from __future__ import annotations
import argparse
import json
import math
import re
import sys
from io import BytesIO
from pathlib import Path
from statistics import median

REPO = Path(__file__).resolve().parent.parent

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ──────────────────────────────────────────────────────────────────────────
# Data loaders + helpers
# ──────────────────────────────────────────────────────────────────────────
def load_data():
    with (REPO / 'ward_data.json').open(encoding='utf-8') as f:
        wd = json.load(f)
    try:
        with (REPO / 'lsoa_data.json').open(encoding='utf-8') as f:
            ld = json.load(f)
    except FileNotFoundError:
        ld = {}
    # Build LSOA → ward_code lookup from lsoa_boundaries.geojson.
    # lsoa_data.json doesn't carry ward_code on each record; the geojson does.
    lsoa_ward = {}
    try:
        with (REPO / 'lsoa_boundaries.geojson').open(encoding='utf-8') as f:
            gj = json.load(f)
        for feat in (gj.get('features') or []):
            p = feat.get('properties') or {}
            code = p.get('LSOA21CD') or p.get('code') or p.get('LSOA11CD')
            wcode = p.get('ward_code')
            if code and wcode:
                lsoa_ward[code] = wcode
    except FileNotFoundError:
        pass
    return wd, ld, lsoa_ward


def find_ward(wd, name_or_code):
    nlow = name_or_code.lower().strip()
    # Try exact code match
    if name_or_code in wd['wards']:
        return name_or_code, wd['wards'][name_or_code]
    # Try name match (case-insensitive)
    for code, w in wd['wards'].items():
        if w.get('name', '').lower() == nlow:
            return code, w
    # Try contains
    for code, w in wd['wards'].items():
        if nlow in w.get('name', '').lower():
            return code, w
    return None, None


def slug(s):
    return re.sub(r'[^a-z0-9]+', '-', (s or '').lower()).strip('-')


def fnum(v):
    if v is None or v == '':
        return None
    try:
        return float(v)
    except (TypeError, ValueError):
        return None


def ward_vals(wd, key):
    """All NWL ward values for indicator key (skipping None)."""
    out = []
    for w in wd['wards'].values():
        v = fnum((w.get('indicators') or {}).get(key))
        if v is not None:
            out.append(v)
    return out


def percentile(vals, x):
    if x is None or not vals:
        return None
    below = sum(1 for v in vals if v < x)
    return round(100 * below / len(vals))


# ──────────────────────────────────────────────────────────────────────────
# Chart renderers (matplotlib → PNG in-memory)
# ──────────────────────────────────────────────────────────────────────────
NAVY = '#1a2640'
INK = '#1a1a1a'
MUTED = '#888888'
GOOD = '#2c8a4a'
BAD = '#c9534a'
NEUTRAL = '#6b7a99'


def _save_png(fig, dpi=180):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


def chart_stacked(parts, title, width=4.6, height=1.2):
    """parts: [(label, value, colour)]"""
    fig, ax = plt.subplots(figsize=(width, height))
    total = sum((v or 0) for _, v, _ in parts) or 1
    x = 0
    for lbl, v, c in parts:
        v = v or 0
        w = v / total
        ax.barh(0, w, left=x, height=0.55, color=c, edgecolor='white', linewidth=0.5)
        if w > 0.07:
            ax.text(x + w/2, 0, f'{v:.0f}%', ha='center', va='center',
                    color='white', fontsize=9, fontweight='bold')
        x += w
    ax.set_xlim(0, 1)
    ax.set_ylim(-0.6, 0.6)
    ax.set_yticks([])
    ax.set_xticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.set_title(title, fontsize=10, color=INK, loc='left', pad=8, fontweight='bold')
    # Legend underneath
    handles = [mpatches.Patch(color=c, label=f'{lbl} {v:.0f}%' if v else lbl)
               for lbl, v, c in parts]
    ax.legend(handles=handles, loc='upper center', bbox_to_anchor=(0.5, -0.15),
              ncol=min(len(parts), 5), fontsize=8, frameon=False)
    return _save_png(fig)


def chart_compare(rows, title, width=4.8, height=2.4, unit='%'):
    """rows: [{l, v, ref, wh}] — wh:True means higher=worse."""
    rows = [r for r in rows if r.get('v') is not None or r.get('ref') is not None]
    if not rows:
        return None
    fig, ax = plt.subplots(figsize=(width, height))
    n = len(rows)
    y_pos = np.arange(n)[::-1]
    vals = [r.get('v') or 0 for r in rows]
    refs = [r.get('ref') or 0 for r in rows]
    dmax = max(max(vals), max(refs), 0.0001) * 1.15
    for i, r in enumerate(rows):
        v = r.get('v')
        ref = r.get('ref')
        wh = r.get('wh', False)
        c = NEUTRAL
        if v is not None and ref is not None:
            above = v > ref
            c = BAD if (wh == above) else GOOD
        ax.barh(y_pos[i], v if v is not None else 0, color=c, height=0.6,
                edgecolor='none')
        if ref is not None:
            ax.plot([ref, ref], [y_pos[i] - 0.35, y_pos[i] + 0.35],
                    color=INK, linewidth=1.5)
        if v is not None:
            ax.text(v + dmax * 0.01, y_pos[i], f'{v:.1f}{unit}',
                    va='center', fontsize=8.5, color=INK)
    ax.set_yticks(y_pos)
    ax.set_yticklabels([r['l'] for r in rows], fontsize=9, color=INK)
    ax.set_xlim(0, dmax)
    ax.set_xticks([])
    for spine in ['top', 'right', 'bottom']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color('#ccc')
    ax.set_title(title, fontsize=10, color=INK, loc='left', pad=8, fontweight='bold')
    ax.tick_params(axis='y', length=0)
    return _save_png(fig)


def chart_hbar(rows, title, width=6.8, height=2.6, unit='', ref_col='ref', sort_desc=True):
    """rows: [{l, v, ref?}]"""
    rows = [r for r in rows if r.get('v') is not None]
    if not rows:
        return None
    if sort_desc:
        rows = sorted(rows, key=lambda r: -(r.get('v') or 0))
    fig, ax = plt.subplots(figsize=(width, height))
    n = len(rows)
    y_pos = np.arange(n)[::-1]
    vals = [r['v'] for r in rows]
    refs = [r.get(ref_col) or 0 for r in rows]
    dmax = max(max(vals), max(refs), 0.0001) * 1.15
    for i, r in enumerate(rows):
        ax.barh(y_pos[i], r['v'], color='#7a3a3a', height=0.6, edgecolor='none')
        if r.get(ref_col) is not None:
            rv = r[ref_col]
            ax.plot([rv, rv], [y_pos[i] - 0.32, y_pos[i] + 0.32],
                    color=INK, linewidth=1.4, linestyle='--')
        ax.text(r['v'] + dmax * 0.01, y_pos[i], f'{r["v"]:,.0f}{unit}',
                va='center', fontsize=8.5, color=INK)
    ax.set_yticks(y_pos)
    ax.set_yticklabels([r['l'] for r in rows], fontsize=9, color=INK)
    ax.set_xlim(0, dmax)
    ax.set_xticks([])
    for spine in ['top', 'right', 'bottom']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color('#ccc')
    ax.set_title(title, fontsize=10, color=INK, loc='left', pad=8, fontweight='bold')
    ax.tick_params(axis='y', length=0)
    return _save_png(fig)


def chart_radar(axes, title, width=4.0, height=4.0):
    """axes: [{l, v, ref, max}]"""
    axes_ok = [a for a in axes if a.get('max')]
    if not axes_ok:
        return None
    n = len(axes_ok)
    angles = np.linspace(0, 2*math.pi, n, endpoint=False).tolist()
    angles += angles[:1]
    vals = [(a.get('v') or 0) / a['max'] for a in axes_ok]
    vals += vals[:1]
    refs = [(a.get('ref') or 0) / a['max'] for a in axes_ok]
    refs += refs[:1]
    fig, ax = plt.subplots(figsize=(width, height), subplot_kw=dict(polar=True))
    # Grid rings
    for r in [0.25, 0.5, 0.75, 1.0]:
        ax.plot([a for a in angles] + [angles[0]],
                [r]*(n+1) + [r],
                color='#e8e8e2', linewidth=0.8)
    ax.fill(angles, vals, color='#b83c3c', alpha=0.30)
    ax.plot(angles, vals, color='#b83c3c', linewidth=1.6)
    # Reference ring
    ax.plot(angles, refs, color=INK, linewidth=1.0, linestyle='--')
    # Axis labels
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels([a['l'] for a in axes_ok], fontsize=8.5, color=INK)
    ax.set_yticks([])
    ax.set_ylim(0, 1.0)
    ax.spines['polar'].set_color('#e8e8e2')
    ax.set_title(title, fontsize=10, color=INK, pad=14, fontweight='bold')
    return _save_png(fig)


# ──────────────────────────────────────────────────────────────────────────
# Slide layouts
# ──────────────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_text_box(slide, x, y, w, h, text, **opts):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(0);  tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    f = run.font
    f.size = Pt(opts.get('size', 14))
    f.bold = opts.get('bold', False)
    rgb = opts.get('color', RGBColor(0x1a, 0x1a, 0x1a))
    if isinstance(rgb, str):
        rgb = RGBColor.from_string(rgb)
    f.color.rgb = rgb
    f.name = opts.get('font', 'Calibri')
    p.alignment = opts.get('align', None)
    return tb


def add_filled_rect(slide, x, y, w, h, fill='ffffff', line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        sh.line.color.rgb = RGBColor.from_string(line)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def slide_cover(prs, ward, wcode, wd, ld):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    inds = ward.get('indicators') or {}
    name = ward.get('name', '?')
    borough = ward.get('lad', '')
    pop = fnum(inds.get('census_population'))
    imd = fnum(inds.get('imd_score'))
    imd_vals = ward_vals(wd, 'imd_score')
    imd_rank = None
    if imd is not None and imd_vals:
        imd_vals_sorted = sorted(imd_vals, reverse=True)
        imd_rank = next((i+1 for i, v in enumerate(imd_vals_sorted) if v <= imd), None)
    n_total = len(imd_vals)
    is_core20 = bool(ward.get('is_core20'))

    # Top navy band
    add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.6), fill=NAVY[1:])
    add_text_box(s, Inches(0.5), Inches(0.18), Inches(11), Inches(0.35),
                 'NW LONDON WARD PROFILE', size=11, bold=True,
                 color=RGBColor(0xca, 0xdc, 0xfc))
    add_text_box(s, Inches(0.5), Inches(0.45), Inches(11), Inches(1.0),
                 name, size=44, bold=True, color=RGBColor(0xff, 0xff, 0xff),
                 font='Calibri')
    sub = borough
    if pop:
        sub += f'   ·   Population {int(pop):,}'
    if imd_rank:
        sub += f'   ·   IMD rank {imd_rank} of {n_total} NWL wards'
    add_text_box(s, Inches(0.5), Inches(1.18), Inches(11), Inches(0.4),
                 sub, size=14, color=RGBColor(0xca, 0xdc, 0xfc))

    # KPI strip — 6 tiles
    kpi_y = Inches(2.0)
    kpi_h = Inches(1.5)
    kpi_w = Inches(2.0)
    kpi_gap = Inches(0.1)
    x = Inches(0.5)
    kpis = [
        ('Population', f'{int(pop):,}' if pop else '—', 'Census 2021'),
        ('IMD score', f'{imd:.1f}' if imd is not None else '—', 'pop-weighted'),
        ('Good health', f'{fnum(inds.get("census_good_health_pct")):.0f}%' if fnum(inds.get('census_good_health_pct')) is not None else '—', 'of residents'),
        ('Bad health', f'{fnum(inds.get("census_bad_health_pct")):.1f}%' if fnum(inds.get('census_bad_health_pct')) is not None else '—', 'of residents'),
        ('Disability', f'{fnum(inds.get("census_disability_any_pct")):.1f}%' if fnum(inds.get('census_disability_any_pct')) is not None else '—', 'any condition'),
        ('Core20 LSOAs', f'{int(fnum(inds.get("core20_lsoa_count")) or 0)}/{int(fnum(inds.get("total_lsoa_count")) or 0)}' if fnum(inds.get('core20_lsoa_count')) is not None else '—', 'most deprived'),
    ]
    for lbl, val, sub in kpis:
        add_filled_rect(s, x, kpi_y, kpi_w, kpi_h, fill='fafaf7', line='eeeeee')
        add_text_box(s, x + Inches(0.1), kpi_y + Inches(0.15), kpi_w - Inches(0.2), Inches(0.25),
                     lbl.upper(), size=10, bold=True, color=RGBColor(0x66, 0x66, 0x66))
        add_text_box(s, x + Inches(0.1), kpi_y + Inches(0.42), kpi_w - Inches(0.2), Inches(0.7),
                     val, size=28, bold=True, color=RGBColor(0x1a, 0x1a, 0x1a))
        add_text_box(s, x + Inches(0.1), kpi_y + Inches(1.10), kpi_w - Inches(0.2), Inches(0.3),
                     sub, size=9, color=RGBColor(0x99, 0x99, 0x99))
        x += kpi_w + kpi_gap

    # Service strip
    svc_y = Inches(3.7)
    svc_h = Inches(1.2)
    svc_w_total = SLIDE_W - Inches(1.0)
    services = [
        ('GP', inds.get('gp_practice_count')),
        ('Pharmacy', inds.get('pharmacy_count')),
        ('Dental', count_layer(ld, wcode, 'dental') if ld else None),  # placeholder
        ('Hospital', None),
        ('Schools', None),  # filled below using --ward filter
        ('Comm ctr', None),
        ('Library', None),
        ('ESOL', None),
    ]
    # Pull schools/CC/lib counts from the JSON files directly
    services[4] = ('Schools', count_by_ward('schools.json', wcode))
    services[5] = ('Comm ctr', count_by_ward('community_centres.json', wcode))
    services[6] = ('Library', count_by_ward('libraries.json', wcode))
    services[7] = ('ESOL', count_esol(wcode))
    nsvc = len(services)
    svc_w = (svc_w_total - Inches(0.1) * (nsvc - 1)) / nsvc
    sx = Inches(0.5)
    for lbl, val in services:
        add_filled_rect(s, sx, svc_y, svc_w, svc_h, fill='fafaf6', line='eeeeee')
        v_str = str(int(val)) if val is not None and val > 0 else '—'
        add_text_box(s, sx, svc_y + Inches(0.18), svc_w, Inches(0.6),
                     v_str, size=24, bold=True,
                     color=RGBColor(0x18, 0x5f, 0xa5) if v_str != '—' else RGBColor(0xbb, 0xbb, 0xbb),
                     align=2)  # 2 = center
        add_text_box(s, sx, svc_y + Inches(0.78), svc_w, Inches(0.3),
                     lbl, size=10, bold=True,
                     color=RGBColor(0x66, 0x66, 0x66), align=2)
        sx += svc_w + Inches(0.1)

    # Footer
    add_text_box(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.4),
                 'Sources: ONS Census 2021 · MHCLG IMD 2019 · MPS via London Datastore · '
                 'NHS Digital ODS · GLA ESOL Planning Map · DfE GIAS · '
                 'London Civic Strength Index R3 · Defra ANGSt · DESNZ LILEE',
                 size=8, color=RGBColor(0x88, 0x88, 0x88))


def count_layer(ld, wcode, layer):
    return None  # placeholder - we don't index dental by ward_code


def count_by_ward(filename, wcode):
    p = REPO / filename
    if not p.exists() or not wcode:
        return None
    try:
        data = json.loads(p.read_text(encoding='utf-8'))
        return sum(1 for r in data if r.get('ward_code') == wcode)
    except Exception:
        return None


def count_esol(wcode):
    # ESOL has no ward_code field; we'd need spatial test, skip for now
    return None


def slide_demographics(prs, ward, wd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    inds = ward.get('indicators') or {}
    add_text_box(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 f'{ward["name"]} — Demographics', size=24, bold=True,
                 color=RGBColor(0x1a, 0x26, 0x40))
    add_text_box(s, Inches(0.5), Inches(0.78), Inches(12), Inches(0.3),
                 'Census 2021 · ward-aggregated', size=11,
                 color=RGBColor(0x88, 0x88, 0x88))

    # Ethnicity stack
    eth = chart_stacked([
        ('White', fnum(inds.get('census_white_pct')), '#5b8db8'),
        ('Asian', fnum(inds.get('census_asian_pct')), '#d47a2c'),
        ('Black', fnum(inds.get('census_black_pct')), '#6b4e71'),
        ('Mixed', fnum(inds.get('census_mixed_pct')), '#8ca757'),
        ('Other', fnum(inds.get('census_other_ethnic_pct')), '#999999'),
    ], 'Ethnicity', width=5.5, height=1.4)
    if eth:
        s.shapes.add_picture(eth, Inches(0.5), Inches(1.4), width=Inches(6))

    # Age stack
    age = chart_stacked([
        ('<16', fnum(inds.get('census_under16_pct')), '#2ea8a0'),
        ('16-64', fnum(inds.get('census_working_age_pct')), '#3b6eaf'),
        ('65+', fnum(inds.get('census_over65_pct')), '#a05a9e'),
    ], 'Age profile', width=5.5, height=1.4)
    if age:
        s.shapes.add_picture(age, Inches(7), Inches(1.4), width=Inches(6))

    # Compare bars
    compare = chart_compare([
        {'l': 'Born outside UK',     'v': fnum(inds.get('census_born_outside_uk_pct')),       'ref': median(ward_vals(wd, 'census_born_outside_uk_pct')) if ward_vals(wd, 'census_born_outside_uk_pct') else None,       'wh': False},
        {'l': 'No qualifications',   'v': fnum(inds.get('census_no_qual_pct')),               'ref': median(ward_vals(wd, 'census_no_qual_pct')) if ward_vals(wd, 'census_no_qual_pct') else None,                 'wh': True},
        {'l': 'Level 4+ qual',       'v': fnum(inds.get('census_level4_qual_pct')),           'ref': median(ward_vals(wd, 'census_level4_qual_pct')) if ward_vals(wd, 'census_level4_qual_pct') else None,             'wh': False},
        {'l': 'Higher managerial',   'v': fnum(inds.get('census_higher_managerial_pct')),     'ref': median(ward_vals(wd, 'census_higher_managerial_pct')) if ward_vals(wd, 'census_higher_managerial_pct') else None, 'wh': False},
        {'l': 'Routine/semi',        'v': fnum(inds.get('census_routine_semi_routine_pct')),  'ref': median(ward_vals(wd, 'census_routine_semi_routine_pct')) if ward_vals(wd, 'census_routine_semi_routine_pct') else None, 'wh': True},
        {'l': 'No-Eng household',    'v': fnum(inds.get('census_english_hh_none_pct')),       'ref': median(ward_vals(wd, 'census_english_hh_none_pct')) if ward_vals(wd, 'census_english_hh_none_pct') else None,         'wh': True},
    ], 'Demographic indicators vs NWL median (red = direction-of-concern, green = strength)',
       width=12, height=3.4)
    if compare:
        s.shapes.add_picture(compare, Inches(0.5), Inches(3.3), width=Inches(12.3))


def slide_health(prs, ward, wd, ld, lsoa_ward, wcode):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    inds = ward.get('indicators') or {}

    add_text_box(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 f'{ward["name"]} — Health & deprivation', size=24, bold=True,
                 color=RGBColor(0x1a, 0x26, 0x40))

    # IMD radar — domains (mean of ward LSOAs).
    # lsoa_data.json doesn't carry ward_code; we use the lsoa_boundaries.geojson
    # lookup built in load_data() to associate each LSOA with its ward.
    domain_keys = ['income_score', 'employment_score', 'education_score',
                   'health_score', 'crime_score', 'barriers_score', 'environment_score']
    domain_labels = ['Income', 'Employment', 'Education', 'Health', 'Crime', 'Barriers', 'Environment']
    ward_doms = {k: 0.0 for k in domain_keys}
    n_lsoa = 0
    # Restrict the NWL pool to NWL LSOAs (those in lsoa_ward) so reference rings
    # are NWL-only rather than London-wide.
    nwl_lsoa_codes = set(lsoa_ward.keys())
    for lcode, ldat in (ld or {}).items():
        if lcode not in nwl_lsoa_codes:
            continue
        d = ldat.get('indicators') if isinstance(ldat, dict) and 'indicators' in ldat else ldat
        if lsoa_ward.get(lcode) != wcode:
            continue
        for k in domain_keys:
            v = fnum(d.get(k))
            if v is not None:
                ward_doms[k] += v
        n_lsoa += 1
    if n_lsoa:
        ward_doms = {k: v/n_lsoa for k, v in ward_doms.items()}

    nwl_doms = {k: [] for k in domain_keys}
    for lcode, ldat in (ld or {}).items():
        if lcode not in nwl_lsoa_codes:
            continue
        d = ldat.get('indicators') if isinstance(ldat, dict) and 'indicators' in ldat else ldat
        for k in domain_keys:
            v = fnum(d.get(k))
            if v is not None:
                nwl_doms[k].append(v)
    nwl_means = {k: (sum(vs)/len(vs) if vs else 0) for k, vs in nwl_doms.items()}
    nwl_max = {k: (max(vs) if vs else 1) for k, vs in nwl_doms.items()}

    radar = chart_radar([
        {'l': lbl, 'v': ward_doms.get(k, 0), 'ref': nwl_means.get(k, 0), 'max': nwl_max.get(k, 1)}
        for lbl, k in zip(domain_labels, domain_keys)
    ], 'IMD domains — mean of ward LSOAs', width=4.5, height=4.5)
    if radar:
        s.shapes.add_picture(radar, Inches(0.4), Inches(1.0), width=Inches(4.8))

    # Health bars
    health = chart_compare([
        {'l': 'Good health',         'v': fnum(inds.get('census_good_health_pct')),           'ref': median(ward_vals(wd, 'census_good_health_pct')) if ward_vals(wd, 'census_good_health_pct') else None,             'wh': False},
        {'l': 'Bad health',          'v': fnum(inds.get('census_bad_health_pct')),            'ref': median(ward_vals(wd, 'census_bad_health_pct')) if ward_vals(wd, 'census_bad_health_pct') else None,              'wh': True},
        {'l': 'Disability (any)',    'v': fnum(inds.get('census_disability_any_pct')),        'ref': median(ward_vals(wd, 'census_disability_any_pct')) if ward_vals(wd, 'census_disability_any_pct') else None,      'wh': True},
        {'l': 'Disability (lot)',    'v': fnum(inds.get('census_disability_lot_pct')),        'ref': median(ward_vals(wd, 'census_disability_lot_pct')) if ward_vals(wd, 'census_disability_lot_pct') else None,      'wh': True},
        {'l': 'Unpaid care',         'v': fnum(inds.get('census_provides_unpaid_care_pct')),  'ref': median(ward_vals(wd, 'census_provides_unpaid_care_pct')) if ward_vals(wd, 'census_provides_unpaid_care_pct') else None, 'wh': False},
        {'l': 'Fuel poverty',        'v': fnum(inds.get('fuel_poverty_pct')),                 'ref': median(ward_vals(wd, 'fuel_poverty_pct')) if ward_vals(wd, 'fuel_poverty_pct') else None,                          'wh': True},
        {'l': 'Unemployed',          'v': fnum(inds.get('census_unemployed_pct')),            'ref': median(ward_vals(wd, 'census_unemployed_pct')) if ward_vals(wd, 'census_unemployed_pct') else None,               'wh': True},
        {'l': 'Claimant rate',       'v': fnum(inds.get('claimant_rate_pct')),                'ref': median(ward_vals(wd, 'claimant_rate_pct')) if ward_vals(wd, 'claimant_rate_pct') else None,                       'wh': True},
    ], 'Health & economic indicators vs NWL median', width=7.5, height=4.5)
    if health:
        s.shapes.add_picture(health, Inches(5.3), Inches(1.0), width=Inches(7.7))


def slide_crime(prs, ward, wd):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    inds = ward.get('indicators') or {}
    add_text_box(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 f'{ward["name"]} — Crime', size=24, bold=True,
                 color=RGBColor(0x1a, 0x26, 0x40))
    crime_total = fnum(inds.get('crime_total'))
    add_text_box(s, Inches(0.5), Inches(0.78), Inches(12), Inches(0.3),
                 f'MPS recorded crime · Apr 2025 to Mar 2026' +
                 (f' · total {int(crime_total):,}' if crime_total else ''),
                 size=11, color=RGBColor(0x88, 0x88, 0x88))

    cats = [
        ('Violence', 'crime_violence_12mo'),
        ('Theft', 'crime_theft_12mo'),
        ('Burglary', 'crime_burglary_12mo'),
        ('Robbery', 'crime_robbery_12mo'),
        ('Vehicle', 'crime_vehicle_12mo'),
        ('Drug offences', 'crime_drug_offences_12mo'),
        ('Public order', 'crime_public_order_12mo'),
        ('Arson & criminal damage', 'crime_arson_12mo'),
        ('Possession of weapons', 'crime_weapons_12mo'),
        ('Sexual offences', 'crime_sexual_offences_12mo'),
        ('Misc. against society', 'crime_misc_society_12mo'),
    ]
    rows = []
    for lbl, k in cats:
        v = fnum(inds.get(k))
        ref = median(ward_vals(wd, k)) if ward_vals(wd, k) else None
        if v is not None:
            rows.append({'l': lbl, 'v': v, 'ref': ref})
    img = chart_hbar(rows, 'By Home Office category — sorted by ward count, dashed = NWL median',
                     width=12, height=4.7, sort_desc=True)
    if img:
        s.shapes.add_picture(img, Inches(0.5), Inches(1.3), width=Inches(12.3))


def slide_civic(prs, ward, wd, ld, lsoa_ward, wcode):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    inds = ward.get('indicators') or {}
    add_text_box(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 f'{ward["name"]} — Civic strength & community', size=24, bold=True,
                 color=RGBColor(0x1a, 0x26, 0x40))

    # Civic — show as percentile in NWL
    civic_fields = [
        ('cst_number_of_community_sport_and_physical_activity_offerings', 'Sport & activity'),
        ('cst_number_of_community_interest_companies_cics', 'Community Interest Companies'),
        ('cst_number_and_proximity_of_libraries', 'Libraries (count & proximity)'),
        ('cst_number_and_proximity_of_community_centres', 'Community centres'),
        ('cst_number_and_proximity_of_cultural_spaces', 'Cultural spaces'),
        ('cst_number_of_faith_centres', 'Faith centres'),
        ('cst_passive_green_space', 'Passive green space'),
    ]
    civic_rows = []
    for k, lbl in civic_fields:
        v = fnum(inds.get(k))
        pct = percentile(ward_vals(wd, k), v)
        if pct is not None:
            civic_rows.append({'l': lbl, 'v': pct, 'ref': 50})
    img = chart_hbar(civic_rows, 'Civic strength — NWL percentile (higher = stronger, dashed = median)',
                     width=7, height=3.6, unit='', ref_col='ref', sort_desc=True)
    if img:
        s.shapes.add_picture(img, Inches(0.4), Inches(1.0), width=Inches(7.0))

    # Green & blue space — mean of ward LSOAs (using LSOA→ward lookup)
    gb_fields = [
        ('gb_commitment_pct', 'Green or blue ≤15min'),
        ('green_commitment_pct', 'Green ≤15min'),
        ('green_doorstep_pct', 'Doorstep (<200m)'),
        ('green_local_pct', 'Local (<300m, ≥2 ha)'),
        ('green_neighbourhood_pct', 'Neighbourhood (<1km)'),
        ('blue_commitment_pct', 'Blue ≤15min'),
    ]
    gb_rows = []
    for key, lbl in gb_fields:
        vs = []
        for lcode, ldat in (ld or {}).items():
            if lsoa_ward.get(lcode) != wcode:
                continue
            d = ldat.get('indicators') if isinstance(ldat, dict) and 'indicators' in ldat else ldat
            v = fnum(d.get(key))
            if v is not None:
                vs.append(v)
        if vs:
            gb_rows.append({'l': lbl, 'v': sum(vs)/len(vs)})
    img2 = chart_hbar(gb_rows, 'Green & blue space access — mean of ward LSOAs (% UPRNs)',
                      width=5, height=3.6, unit='%', ref_col='__no__', sort_desc=False)
    if img2:
        s.shapes.add_picture(img2, Inches(7.7), Inches(1.0), width=Inches(5.3))


# ──────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────
def build_one(ward_arg, out_dir):
    wd, ld, lsoa_ward = load_data()
    code, ward = find_ward(wd, ward_arg)
    if not ward:
        print(f"No ward matching {ward_arg!r}", file=sys.stderr)
        return False
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f'nwl-ward-{slug(ward["name"])}.pptx'

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, ward, code, wd, ld)
    slide_demographics(prs, ward, wd)
    slide_health(prs, ward, wd, ld, lsoa_ward, code)
    slide_crime(prs, ward, wd)
    slide_civic(prs, ward, wd, ld, lsoa_ward, code)

    prs.save(str(out))
    print(f"  wrote {out}")
    return True


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                  formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('--ward', help='ward name or WD24CD (e.g. "St James'+chr(39)+'s")')
    ap.add_argument('--all', action='store_true', help='build a deck for every NWL ward')
    ap.add_argument('--list', action='store_true', help='list available wards and exit')
    ap.add_argument('--out', default='output/ward_decks',
                    help='output directory (relative to repo root)')
    args = ap.parse_args()

    out_dir = (REPO / args.out).resolve()

    if args.list:
        wd, _, _ = load_data()
        for w in sorted(wd['wards'].values(), key=lambda x: (x.get('lad', ''), x.get('name', ''))):
            print(f"  [{w.get('lad','?'):28}] {w.get('name')}")
        return 0

    if args.all:
        wd, _, _ = load_data()
        n = len(wd['wards'])
        ok = 0
        for i, (code, w) in enumerate(wd['wards'].items(), 1):
            print(f"[{i:3}/{n}] {w.get('name')} ({code})")
            if build_one(code, out_dir):
                ok += 1
        print(f"Built {ok}/{n} decks in {out_dir}")
        return 0

    if not args.ward:
        ap.error('one of --ward, --all, or --list is required')

    return 0 if build_one(args.ward, out_dir) else 1


if __name__ == '__main__':
    sys.exit(main())
